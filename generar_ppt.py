import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = r'C:\Users\dmoov\Documents\Programas\costos-gastronomia\Culinary_Sistema_Costos_Presentacion.pptx'

# ── Colores ──────────────────────────────────────────────────────
C_BG     = RGBColor(0x14, 0x14, 0x14)   # fondo oscuro
C_BG2    = RGBColor(0x1e, 0x1e, 0x1e)   # fondo card
C_GOLD   = RGBColor(0xfb, 0xbf, 0x24)   # dorado
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_MUTED  = RGBColor(0x9c, 0xa3, 0xaf)   # gris claro
C_RED    = RGBColor(0xef, 0x44, 0x44)
C_GREEN  = RGBColor(0x22, 0xc5, 0x5e)
C_BLUE   = RGBColor(0x38, 0xbd, 0xf8)
C_ACCENT = RGBColor(0x2d, 0x2d, 0x2d)   # fondo sección

# ── Slide dimensions: widescreen 16:9 ────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

blank_layout = prs.slide_layouts[6]  # completamente en blanco

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, x, y, w, h, fill=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def txbox(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color or C_WHITE
    return tb

def bg(slide, color=None):
    bg_shape = slide.shapes.add_shape(1, 0, 0, W, H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color or C_BG
    bg_shape.line.fill.background()
    return bg_shape

def hline(slide, x, y, w, color=None, thickness=0.03):
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color or C_GOLD
    line.line.fill.background()

def card(slide, x, y, w, h, title='', value='', sub='', color=None, vcolor=None):
    rect(slide, x, y, w, h, fill=C_ACCENT)
    hline(slide, x, y, w, color=color or C_GOLD, thickness=0.04)
    if title:
        txbox(slide, title, x+0.1, y+0.12, w-0.2, 0.35, size=10, color=C_MUTED, align=PP_ALIGN.LEFT)
    if value:
        txbox(slide, value, x+0.1, y+0.42, w-0.2, 0.55, size=22, bold=True, color=vcolor or C_GOLD, align=PP_ALIGN.LEFT)
    if sub:
        txbox(slide, sub, x+0.1, y+0.92, w-0.2, 0.3, size=9, color=C_MUTED, align=PP_ALIGN.LEFT)

def slide_number(slide, n, total):
    txbox(slide, f'{n} / {total}', 12.2, 7.1, 1.0, 0.3, size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)

def footer(slide, text='Instituto Culinary · Sistema de Costos Gastronómicos · 2026'):
    hline(slide, 0.3, 7.15, 12.73, color=RGBColor(0x33,0x33,0x33))
    txbox(slide, text, 0.4, 7.2, 10, 0.25, size=8, color=C_MUTED)

TOTAL_SLIDES = 13

# ═══════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
# Banda dorada vertical izquierda
rect(sl, 0, 0, 0.08, 7.5, fill=C_GOLD)
# Franja semi-oscura derecha
rect(sl, 0.08, 0, 13.25, 7.5, fill=RGBColor(0x16,0x16,0x16))
# Patrón diagonal decorativo (rectángulos dorados sutiles)
for i, (xi, yi) in enumerate([(10.5,0.2),(11.0,1.1),(11.5,0.6),(12.0,1.5),(12.5,0.9),(11.8,2.3),(10.8,2.8)]):
    r = sl.shapes.add_shape(1, Inches(xi), Inches(yi), Inches(0.08), Inches(0.8))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xfb,0xbf,0x24)
    r.line.fill.background()
    from pptx.util import Pt as uPt
    r.fill.fore_color.theme_color  # just access to avoid warning

# Etiqueta institución
txbox(sl, 'INSTITUTO CULINARY', 0.5, 1.3, 10, 0.5, size=13, bold=False, color=C_MUTED, font='Calibri')
hline(sl, 0.5, 1.85, 5.5)
# Título principal
txbox(sl, 'SISTEMA DE COSTOS', 0.5, 2.0, 12, 0.9, size=46, bold=True, color=C_WHITE, font='Calibri')
txbox(sl, 'GASTRONÓMICOS', 0.5, 2.85, 12, 0.9, size=46, bold=True, color=C_GOLD, font='Calibri')
# Subtítulo
txbox(sl, 'Gestión de costos de materia prima · Sede Santiago · Período 2026', 0.5, 3.85, 11, 0.4, size=14, color=C_MUTED, font='Calibri')
# Línea separadora
hline(sl, 0.5, 4.4, 5.5)
# Metadatos
txbox(sl, 'Junio 2026', 0.5, 4.6, 4, 0.35, size=11, color=C_MUTED)
txbox(sl, 'Versión 1.0 — Uso Interno', 0.5, 4.95, 6, 0.35, size=11, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 2, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, 'CONTENIDO', 0.5, 0.3, 5, 0.5, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Agenda de la presentación', 0.5, 0.65, 8, 0.4, size=18, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

items = [
    ('01', 'Contexto y necesidad'),
    ('02', 'La solución: el sistema'),
    ('03', 'Arquitectura técnica'),
    ('04', 'Datos del período 2026'),
    ('05', 'Dashboard y KPIs principales'),
    ('06', 'Análisis de costos por categoría'),
    ('07', 'Costos por tipo de taller'),
    ('08', 'Hallazgos y observaciones'),
    ('09', 'Próximos pasos'),
]
col1 = items[:5]
col2 = items[5:]
for i, (num, txt) in enumerate(col1):
    y = 1.4 + i * 0.95
    txbox(sl, num, 0.5, y, 0.5, 0.4, size=20, bold=True, color=C_GOLD)
    txbox(sl, txt, 1.1, y+0.05, 5.5, 0.4, size=14, color=C_WHITE)
    hline(sl, 0.5, y+0.72, 5.8, color=RGBColor(0x2e,0x2e,0x2e))

for i, (num, txt) in enumerate(col2):
    y = 1.4 + i * 0.95
    txbox(sl, num, 6.9, y, 0.5, 0.4, size=20, bold=True, color=C_GOLD)
    txbox(sl, txt, 7.5, y+0.05, 5.5, 0.4, size=14, color=C_WHITE)
    hline(sl, 6.9, y+0.72, 6.1, color=RGBColor(0x2e,0x2e,0x2e))

# ═══════════════════════════════════════════
# SLIDE 3 — CONTEXTO Y NECESIDAD
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 3, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '01 · CONTEXTO', 0.5, 0.3, 5, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Problema y necesidad', 0.5, 0.65, 10, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

problemas = [
    ('Falta de visibilidad', 'Los costos de materia prima se manejaban en planillas Excel sin integración ni historial.'),
    ('Cálculos manuales', 'Cada taller calculaba sus costos de forma independiente, sin una metodología unificada.'),
    ('Sin trazabilidad', 'Era imposible comparar costos entre talleres, secciones o períodos de tiempo.'),
    ('Toma de decisiones', 'La falta de datos precisos dificultaba la planificación presupuestaria y la gestión académica.'),
]
for i, (titulo, desc) in enumerate(problemas):
    x = 0.5 + (i % 2) * 6.4
    y = 1.5 + (i // 2) * 2.6
    rect(sl, x, y, 6.0, 2.3, fill=C_ACCENT)
    r = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.05), Inches(2.3))
    r.fill.solid(); r.fill.fore_color.rgb = C_GOLD; r.line.fill.background()
    txbox(sl, titulo, x+0.2, y+0.15, 5.5, 0.4, size=13, bold=True, color=C_GOLD)
    txbox(sl, desc, x+0.2, y+0.6, 5.5, 1.5, size=11, color=C_WHITE)

# ═══════════════════════════════════════════
# SLIDE 4 — LA SOLUCIÓN
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 4, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '02 · LA SOLUCIÓN', 0.5, 0.3, 5, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Sistema de Costos Gastronómicos', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

txbox(sl, 'Una plataforma web centralizada que permite gestionar, visualizar y analizar los costos de materia prima\nde todos los talleres culinarios del Instituto, en tiempo real y con datos históricos.', 0.5, 1.4, 12.5, 0.8, size=13, color=C_MUTED)

funciones = [
    ('📊', 'Dashboard\nInteractivo', 'KPIs, gráficos y filtros por taller, mes y semana'),
    ('🧂', 'Maestro de\nIngredientes', '613 ingredientes con precio, unidad y categoría'),
    ('📋', 'Gestión de\nOPs', 'Registro de recetas por orden de producto con costo automático'),
    ('📚', 'Carreras y\nSemestres', 'Estructura académica con asignación de talleres'),
    ('🛒', 'Lista de\nCompras', 'Alertas de stock bajo y necesidades de reposición'),
    ('📈', 'Informes\nExportables', 'PDF, Excel, plantillas de importación masiva'),
]
for i, (icon, titulo, desc) in enumerate(funciones):
    x = 0.5 + (i % 3) * 4.2
    y = 2.5 + (i // 3) * 2.3
    rect(sl, x, y, 3.9, 2.0, fill=C_ACCENT)
    txbox(sl, icon, x+0.15, y+0.15, 0.6, 0.5, size=22)
    txbox(sl, titulo, x+0.85, y+0.12, 2.8, 0.6, size=12, bold=True, color=C_WHITE)
    txbox(sl, desc, x+0.15, y+0.9, 3.6, 0.9, size=10, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 5 — ARQUITECTURA TÉCNICA
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 5, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '03 · ARQUITECTURA', 0.5, 0.3, 5, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Stack tecnológico', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

capas = [
    ('FRONTEND', 'index.html (Single Page App)', '• Vanilla JS + Chart.js + SheetJS\n• Sin frameworks, cero dependencias externas\n• Desplegado en PythonAnywhere', C_BLUE),
    ('BACKEND / API', 'Supabase PostgREST', '• REST API automática desde PostgreSQL\n• Autenticación con anon key\n• Vistas SQL para cálculos complejos', C_GREEN),
    ('BASE DE DATOS', 'PostgreSQL (Supabase)', '• 12 tablas + 6 vistas calculadas\n• Triggers para historial de precios\n• Funciones para costeo automático', C_GOLD),
]
for i, (capa, tech, desc, color) in enumerate(capas):
    x = 0.5 + i * 4.25
    rect(sl, x, 1.5, 4.0, 5.5, fill=C_ACCENT)
    hline(sl, x, 1.5, 4.0, color=color, thickness=0.05)
    txbox(sl, capa, x+0.2, 1.65, 3.6, 0.35, size=10, bold=True, color=color)
    txbox(sl, tech, x+0.2, 2.05, 3.6, 0.45, size=13, bold=True, color=C_WHITE)
    hline(sl, x+0.2, 2.55, 3.4, color=RGBColor(0x33,0x33,0x33))
    txbox(sl, desc, x+0.2, 2.75, 3.6, 2.5, size=11, color=C_MUTED)

txbox(sl, '→', 4.52, 3.9, 0.3, 0.4, size=18, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
txbox(sl, '→', 8.77, 3.9, 0.3, 0.4, size=18, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
txbox(sl, '💻 Navegador del usuario', 0.5, 7.05, 4.0, 0.3, size=9, color=C_MUTED, align=PP_ALIGN.CENTER)
txbox(sl, '☁️ Cloud API', 4.75, 7.05, 4.0, 0.3, size=9, color=C_MUTED, align=PP_ALIGN.CENTER)
txbox(sl, '🗄️ Cloud Database', 9.0, 7.05, 4.0, 0.3, size=9, color=C_MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 6 — DATOS DEL PERÍODO
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 6, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '04 · DATOS DEL PERÍODO', 0.5, 0.3, 8, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Carga real — Sede Santiago 2026', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

kpis = [
    ('613', 'Ingredientes', 'en maestro con precio', C_GOLD),
    ('22', 'Proveedores', 'activos registrados', C_BLUE),
    ('14', 'Tipos de Taller', 'especializaciones', C_GREEN),
    ('179', 'OPs', 'órdenes de producto únicas', C_GOLD),
    ('23.213', 'Registros', 'ingrediente × sesión', C_BLUE),
    ('44', 'Semanas', 'calendario 2026', C_GREEN),
]
cw = 4.0
for i, (val, tit, sub, color) in enumerate(kpis):
    x = 0.5 + (i % 3) * (cw + 0.25)
    y = 1.5 + (i // 3) * 2.6
    card(sl, x, y, cw, 2.3, title=tit, value=val, sub=sub, color=color, vcolor=color)

# ═══════════════════════════════════════════
# SLIDE 7 — DASHBOARD / KPIs
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 7, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '05 · DASHBOARD', 0.5, 0.3, 5, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'KPIs principales del sistema', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

# KPI grande izquierda
rect(sl, 0.5, 1.5, 6.2, 3.0, fill=C_ACCENT)
hline(sl, 0.5, 1.5, 6.2, color=C_GOLD, thickness=0.07)
txbox(sl, 'GASTO TOTAL MATERIA PRIMA', 0.7, 1.7, 5.8, 0.4, size=11, bold=True, color=C_MUTED)
txbox(sl, '$7.075.872', 0.7, 2.15, 5.8, 1.1, size=44, bold=True, color=C_GOLD)
txbox(sl, 'Período marzo → diciembre 2026', 0.7, 3.2, 5.8, 0.35, size=11, color=C_MUTED)

# KPIs secundarios derecha
right_kpis = [
    ('$328', 'Costo promedio por alumno / sesión', C_BLUE),
    ('1.067', 'Total alumnos matriculados', C_GREEN),
    ('9 de 14', 'Talleres con recetas cargadas', C_GOLD),
]
for i, (val, lbl, color) in enumerate(right_kpis):
    y = 1.5 + i * 1.55
    rect(sl, 7.0, y, 5.8, 1.4, fill=C_ACCENT)
    hline(sl, 7.0, y, 5.8, color=color, thickness=0.04)
    txbox(sl, val, 7.2, y+0.2, 3.0, 0.65, size=28, bold=True, color=color)
    txbox(sl, lbl, 7.2, y+0.85, 5.3, 0.35, size=11, color=C_MUTED)

# Distribución mensual
rect(sl, 0.5, 4.7, 12.5, 2.5, fill=C_ACCENT)
txbox(sl, 'DISTRIBUCIÓN MENSUAL DE GASTO ($)', 0.7, 4.85, 12, 0.3, size=10, bold=True, color=C_MUTED)
meses = [
    ('Mar', 1014475), ('Abr', 965449), ('May', 1178684), ('Jun', 1076417),
    ('Jul', 348671), ('Ago', 528122), ('Sep', 593151), ('Oct', 458301),
    ('Nov', 588237), ('Dic', 324365),
]
max_val = max(v for _,v in meses)
bar_area_w = 11.5
bar_w = bar_area_w / len(meses) - 0.05
for i, (mes, val) in enumerate(meses):
    h_bar = (val / max_val) * 1.3
    x_bar = 0.8 + i * (bar_area_w / len(meses))
    y_bar = 6.85 - h_bar
    bshape = sl.shapes.add_shape(1, Inches(x_bar), Inches(y_bar), Inches(bar_w), Inches(h_bar))
    bshape.fill.solid()
    bshape.fill.fore_color.rgb = C_GOLD if val == max_val else RGBColor(0x78,0x50,0x10)
    bshape.line.fill.background()
    txbox(sl, mes, x_bar, 6.9, bar_w+0.05, 0.2, size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 8 — CATEGORÍAS
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 8, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '06 · ANÁLISIS DE COSTOS', 0.5, 0.3, 8, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Gasto por categoría de ingrediente', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

total = 7075872
categorias = [
    ('ABARROTE',        2040961, '#f59e0b'),
    ('LACTEO',          1241601, '#3b82f6'),
    ('CHOCOLATE',        824444, '#8b5cf6'),
    ('FRUTA VERDURA',    705926, '#22c55e'),
    ('PESCADO MARISCO',  389899, '#06b6d4'),
    ('VACUNO',           346539, '#ef4444'),
    ('HUEVO',            332590, '#f97316'),
    ('CONGELADO',        303393, '#a855f7'),
    ('ART. COCINA',      273737, '#6366f1'),
    ('RESTO (25 cat.)', 617278, '#6b7280'),
]
for i, (cat, monto, hex_color) in enumerate(categorias):
    y = 1.4 + i * 0.57
    pct = monto / total * 100
    bar_w = pct / 100 * 7.5
    r = RGBColor(int(hex_color[1:3],16), int(hex_color[3:5],16), int(hex_color[5:7],16))
    bg_bar = sl.shapes.add_shape(1, Inches(4.5), Inches(y+0.08), Inches(7.6), Inches(0.35))
    bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = C_ACCENT; bg_bar.line.fill.background()
    fill_bar = sl.shapes.add_shape(1, Inches(4.5), Inches(y+0.08), Inches(max(bar_w,0.05)), Inches(0.35))
    fill_bar.fill.solid(); fill_bar.fill.fore_color.rgb = r; fill_bar.line.fill.background()
    txbox(sl, cat, 0.5, y, 3.8, 0.4, size=11, color=C_WHITE)
    txbox(sl, f'${monto:,.0f}'.replace(',','.'), 12.2, y, 1.0, 0.4, size=11, color=C_GOLD, align=PP_ALIGN.RIGHT)
    txbox(sl, f'{pct:.1f}%', 4.55, y, 0.8, 0.4, size=9, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 9 — TALLERES
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 9, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '07 · ANÁLISIS POR TALLER', 0.5, 0.3, 8, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Costo total y costo por alumno según tipo de taller', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

talleres_data = [
    ('TIP',   'Introductorio Pastelería',    1119984, 471,  140),
    ('TBD',   'Banquetería Dulce',           1062113, 909,   64),
    ('TCCHO', 'Confitería y Chocolatería',    975967, 744,   64),
    ('TCCM',  'Cocina Chilena y del Mundo',   901539, 440,  100),
    ('TTPIP', 'Técnicas Innovadoras Past.',   752045,1106,   40),
    ('TIC',   'Introductorio Cocina',         723792, 275,  155),
    ('TBS',   'Banquetería Salada',           657917, 376,   96),
    ('TTPIC', 'Técnicas Innovadoras Cocina',  486289, 715,   40),
    ('TPE',   'Panadería Especialidad',       396225, 254,   76),
]
max_costo = max(t[2] for t in talleres_data)
txbox(sl, 'TALLER', 0.4, 1.4, 1.2, 0.3, size=9, bold=True, color=C_MUTED)
txbox(sl, 'COSTO TOTAL', 1.65, 1.4, 5.5, 0.3, size=9, bold=True, color=C_MUTED)
txbox(sl, 'SESIONES', 9.5, 1.4, 1.2, 0.3, size=9, bold=True, color=C_MUTED)
txbox(sl, 'CPA', 11.2, 1.4, 1.9, 0.3, size=9, bold=True, color=C_MUTED, align=PP_ALIGN.RIGHT)

for i, (cod, nom, costo, cpa, ses) in enumerate(talleres_data):
    y = 1.8 + i * 0.58
    fill_w = (costo / max_costo) * 7.5
    # barra fondo
    b = sl.shapes.add_shape(1, Inches(1.6), Inches(y+0.12), Inches(7.5), Inches(0.3))
    b.fill.solid(); b.fill.fore_color.rgb = C_ACCENT; b.line.fill.background()
    # barra valor
    bv = sl.shapes.add_shape(1, Inches(1.6), Inches(y+0.12), Inches(fill_w), Inches(0.3))
    bv.fill.solid()
    bv.fill.fore_color.rgb = C_GOLD if i == 0 else RGBColor(0x78,0x50,0x10)
    bv.line.fill.background()
    txbox(sl, cod, 0.4, y, 1.1, 0.45, size=10, bold=True, color=C_GOLD)
    txbox(sl, nom, 0.4, y+0.28, 1.1, 0.25, size=7, color=C_MUTED)
    txbox(sl, f'${costo:,.0f}'.replace(',','.'), 1.65, y, 3.0, 0.4, size=10, color=C_WHITE)
    txbox(sl, str(ses), 9.5, y, 0.8, 0.4, size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
    txbox(sl, f'${cpa:,.0f}'.replace(',','.'), 11.2, y, 1.9, 0.4, size=11, bold=True,
          color=C_RED if cpa > 900 else C_GREEN if cpa < 400 else C_GOLD, align=PP_ALIGN.RIGHT)

hline(sl, 0.4, 7.0, 12.5, color=RGBColor(0x33,0x33,0x33))
txbox(sl, '🟢 CPA < $400  🟡 CPA $400–$900  🔴 CPA > $900', 0.4, 7.05, 9, 0.3, size=9, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 10 — TOP SESIONES
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 10, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '07 · ANÁLISIS POR TALLER', 0.5, 0.3, 8, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Top 5 sesiones con mayor costo total', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

top5 = [
    ('🥇', 'TTPICSES3', 'TTPIC', '2026-05-05', 16, 27846, 1740),
    ('🥈', 'TTPICSES3', 'TTPIC', '2026-03-10', 16, 27846, 1740),
    ('🥉', 'TTPICSES3', 'TTPIC', '2026-05-05', 19, 27846, 1466),
    ('4°', 'TTPICSES3', 'TTPIC', '2026-03-10', 19, 27846, 1466),
    ('5°', 'TBSSES12',  'TBS',  '2026-07-07', 18, 23311, 1295),
]
headers = ['', 'Código', 'Taller', 'Fecha', 'Alumnos', 'Costo Total', 'CPA']
col_x   = [0.4, 0.9, 2.5, 4.5, 6.5, 8.3, 11.2]
col_w   = [0.4, 1.5, 1.9, 1.9, 1.7, 2.8, 1.8]
for j, h in enumerate(headers):
    txbox(sl, h, col_x[j], 1.5, col_w[j], 0.35, size=10, bold=True, color=C_MUTED)
hline(sl, 0.4, 1.9, 12.5)

medals_color = [C_GOLD, C_MUTED, RGBColor(0xcd,0x7f,0x32), C_MUTED, C_MUTED]
for i, (medal, cod, tall, fecha, alum, costo, cpa) in enumerate(top5):
    y = 2.1 + i * 0.95
    if i % 2 == 0:
        bg_row = sl.shapes.add_shape(1, Inches(0.4), Inches(y-0.05), Inches(12.5), Inches(0.88))
        bg_row.fill.solid(); bg_row.fill.fore_color.rgb = C_ACCENT; bg_row.line.fill.background()
    txbox(sl, medal, col_x[0], y, col_w[0], 0.6, size=18, align=PP_ALIGN.CENTER)
    txbox(sl, cod, col_x[1], y, col_w[1], 0.6, size=11, color=C_WHITE)
    txbox(sl, tall, col_x[2], y, col_w[2], 0.6, size=11, color=medals_color[i], bold=True)
    txbox(sl, fecha, col_x[3], y, col_w[3], 0.6, size=11, color=C_MUTED)
    txbox(sl, str(alum), col_x[4], y, col_w[4], 0.6, size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, f'${costo:,.0f}'.replace(',','.'), col_x[5], y, col_w[5], 0.6, size=14, bold=True, color=C_RED)
    txbox(sl, f'${cpa:,.0f}'.replace(',','.'), col_x[6], y, col_w[6], 0.6, size=14, bold=True, color=C_GOLD)

txbox(sl, 'Nota: TTPIC (Técnicas y Procesos Innovadores en Cocina) concentra las sesiones de mayor costo por su uso de ingredientes especializados.', 0.4, 6.9, 12.5, 0.4, size=9, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 11 — TOP INGREDIENTES
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 11, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '06 · ANÁLISIS DE COSTOS', 0.5, 0.3, 8, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Top 10 ingredientes por gasto acumulado', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

top_ing = [
    ('MANTEQUILLA SIN SAL', 'LACTEO',    341058, 'kg'),
    ('CHOCOLATE',           'CHOCOLATE', 270822, 'kg'),
    ('CREMA',               'LACTEO',    254698, 'kg'),
    ('HUEVO',               'HUEVO',     251073, 'un'),
    ('LECHE',               'LACTEO',    226814, 'L'),
    ('HARINA',              'ABARROTE',  195432, 'kg'),
    ('AZUCAR',              'ABARROTE',  187341, 'kg'),
    ('VACUNO LOMO',         'VACUNO',    162890, 'kg'),
    ('SALMON',              'PESCADO',   143201, 'kg'),
    ('COUVERTURE OSCURO',   'CHOCOLATE', 138765, 'kg'),
]
max_ing = top_ing[0][2]
for i, (nom, cat, monto, unid) in enumerate(top_ing):
    y = 1.5 + i * 0.55
    fill_w = (monto / max_ing) * 7.0
    b = sl.shapes.add_shape(1, Inches(4.6), Inches(y+0.1), Inches(7.1), Inches(0.3))
    b.fill.solid(); b.fill.fore_color.rgb = C_ACCENT; b.line.fill.background()
    bv = sl.shapes.add_shape(1, Inches(4.6), Inches(y+0.1), Inches(max(fill_w,0.05)), Inches(0.3))
    bv.fill.solid(); bv.fill.fore_color.rgb = C_GOLD if i == 0 else RGBColor(0x78,0x50,0x10); bv.line.fill.background()
    txbox(sl, f'{i+1:2d}. {nom}', 0.4, y, 4.0, 0.4, size=11, color=C_WHITE)
    txbox(sl, cat, 0.4, y+0.3, 4.0, 0.2, size=8, color=C_MUTED)
    txbox(sl, f'${monto:,.0f}'.replace(',','.'), 11.7, y, 1.4, 0.4, size=11, color=C_GOLD, align=PP_ALIGN.RIGHT)

txbox(sl, '* Montos basados en precios convertidos a unidad base (kg/L/un)', 0.4, 7.1, 10, 0.25, size=8, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 12 — HALLAZGOS
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
footer(sl)
slide_number(sl, 12, TOTAL_SLIDES)
rect(sl, 0, 0, 13.333, 1.2, fill=C_ACCENT)
txbox(sl, '08 · HALLAZGOS', 0.5, 0.3, 5, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Observaciones y conclusiones del análisis', 0.5, 0.65, 11, 0.4, size=22, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.2, 5)

hallazgos = [
    ('✅', C_GREEN,  'Costo promedio razonable',
     'El CPA promedio de $328 por sesión es consistente con los estándares de formación culinaria. El modelo de costeo "por sección" (no por alumno) coincide con el criterio del Excel original.'),
    ('⚠️', C_GOLD,   'TTPIP tiene el mayor costo por alumno',
     'El Taller de Técnicas y Procesos Innovadores en Pastelería tiene el CPA más alto ($1.106), seguido por TBD ($909) y TCCHO ($744). Estos talleres requieren insumos premium.'),
    ('📊', C_BLUE,   'Mayo es el mes con mayor gasto',
     'Con $1.178.684, mayo concentra el mayor gasto de materia prima. Julio registra el menor ($348.671), posiblemente por receso académico o menor número de sesiones.'),
    ('🔧', C_MUTED,  '5 talleres sin recetas cargadas',
     'TMA, TAS, TIR, TMI y TEAC tienen $0 en gasto, lo que indica que sus recetas aún no han sido registradas en el sistema. Se recomienda completar la carga para una visión integral.'),
]
for i, (icon, color, titulo, desc) in enumerate(hallazgos):
    x = 0.5 + (i % 2) * 6.4
    y = 1.5 + (i // 2) * 2.6
    rect(sl, x, y, 6.0, 2.4, fill=C_ACCENT)
    hline(sl, x, y, 6.0, color=color, thickness=0.04)
    txbox(sl, icon + '  ' + titulo, x+0.2, y+0.15, 5.6, 0.45, size=13, bold=True, color=color)
    txbox(sl, desc, x+0.2, y+0.68, 5.6, 1.55, size=11, color=C_MUTED)

# ═══════════════════════════════════════════
# SLIDE 13 — PRÓXIMOS PASOS / CIERRE
# ═══════════════════════════════════════════
sl = add_slide()
bg(sl)
rect(sl, 0, 0, 0.08, 7.5, fill=C_GOLD)
txbox(sl, '09 · PRÓXIMOS PASOS', 0.5, 0.4, 8, 0.4, size=11, color=C_GOLD, bold=True)
txbox(sl, 'Roadmap y mejoras planificadas', 0.5, 0.75, 11, 0.5, size=24, bold=True, color=C_WHITE)
hline(sl, 0.5, 1.35, 6)

pasos = [
    ('CORTO PLAZO', [
        'Completar recetas de TMA, TAS, TIR, TMI y TEAC',
        'Vincular sesiones con semanas del calendario',
        'Actualizar precios de ingredientes para S2 2026',
    ], C_GOLD),
    ('MEDIANO PLAZO', [
        'Módulo de presupuesto vs. real por carrera',
        'Integración con sistema de inventario físico',
        'Exportación automática de informes mensuales',
    ], C_BLUE),
    ('LARGO PLAZO', [
        'Extensión a otras sedes del Instituto',
        'Dashboard comparativo multi-sede',
        'API de integración con sistema administrativo',
    ], C_GREEN),
]
for i, (plazo, items, color) in enumerate(pasos):
    x = 0.5 + i * 4.2
    rect(sl, x, 1.6, 3.9, 5.3, fill=C_ACCENT)
    hline(sl, x, 1.6, 3.9, color=color, thickness=0.06)
    txbox(sl, plazo, x+0.2, 1.75, 3.5, 0.4, size=11, bold=True, color=color)
    for j, item in enumerate(items):
        yj = 2.3 + j * 0.95
        dot = sl.shapes.add_shape(1, Inches(x+0.2), Inches(yj+0.13), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        txbox(sl, item, x+0.45, yj, 3.3, 0.8, size=11, color=C_WHITE)

# Cierre
rect(sl, 0.5, 7.0, 12.5, 0.01, fill=RGBColor(0x33,0x33,0x33))
txbox(sl, 'Instituto Culinary · Sistema de Costos Gastronómicos · Sede Santiago 2026', 0.5, 7.1, 10, 0.3, size=9, color=C_MUTED)
txbox(sl, 'culinary.pythonanywhere.com', 10.0, 7.1, 3.0, 0.3, size=9, color=C_GOLD, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(f'PPT guardado: {OUT}')
